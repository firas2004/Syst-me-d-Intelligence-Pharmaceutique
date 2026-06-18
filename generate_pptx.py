from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Brand Colors ──────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x11, 0x17)   # #0D1117
RED     = RGBColor(0xC0, 0x27, 0x2D)   # #C0272D
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xA0, 0xA0, 0xA0)
DARKCARD= RGBColor(0x16, 0x1B, 0x22)   # #161B22

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── Helpers ───────────────────────────────────────────────────
def add_slide():
    s = prs.slides.add_slide(blank_layout)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def red_bar(slide, y=Inches(0.55), w=Inches(1.2)):
    """Thin red accent line under title."""
    line = slide.shapes.add_shape(1, Inches(0.5), y, w, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

def card(slide, x, y, w, h, title, lines, title_color=RED):
    """Dark card with title + bullet lines."""
    box = slide.shapes.add_shape(1, x, y, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = DARKCARD
    box.line.color.rgb = RED

    # title
    txbox(slide, title, x+Inches(0.12), y+Inches(0.1),
          w-Inches(0.24), Inches(0.45),
          size=13, bold=True, color=title_color)
    # bullets
    tb = slide.shapes.add_textbox(x+Inches(0.12), y+Inches(0.5),
                                   w-Inches(0.24), h-Inches(0.6))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = GREY

def slide_number(slide, n):
    txbox(slide, f"{n}/9", Inches(12.6), Inches(7.1),
          Inches(0.6), Inches(0.3), size=9, color=GREY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — L'Accroche
# ═══════════════════════════════════════════════════════════════
s1 = add_slide()

# Red diagonal accent shape (top-left)
tri = s1.shapes.add_shape(6, Inches(0), Inches(0), Inches(1.5), Inches(7.5))
tri.fill.solid(); tri.fill.fore_color.rgb = RGBColor(0x18, 0x04, 0x04)
tri.line.fill.background()

txbox(s1, "PHARMA", Inches(1.8), Inches(1.8), Inches(6), Inches(1.2),
      size=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
txbox(s1, "INTEL", Inches(5.55), Inches(1.8), Inches(5), Inches(1.2),
      size=60, bold=True, color=RED, align=PP_ALIGN.LEFT)
txbox(s1, "v3.0  ·  Opalia Recordati", Inches(1.8), Inches(2.95),
      Inches(8), Inches(0.5), size=14, color=RED, align=PP_ALIGN.LEFT)

# Red separator line
sep = s1.shapes.add_shape(1, Inches(1.8), Inches(3.55), Inches(3), Pt(2))
sep.fill.solid(); sep.fill.fore_color.rgb = RED; sep.line.fill.background()

txbox(s1, "« L'intelligence au service de la molécule. »",
      Inches(1.8), Inches(3.8), Inches(9), Inches(0.7),
      size=18, italic=True, color=GREY, align=PP_ALIGN.LEFT)

txbox(s1, "13 000+", Inches(1.8), Inches(4.7), Inches(3), Inches(0.9),
      size=40, bold=True, color=RED, align=PP_ALIGN.LEFT)
txbox(s1, "molécules analysées", Inches(1.8), Inches(5.5), Inches(5), Inches(0.4),
      size=12, color=GREY)

txbox(s1, "© 2025 Opalia Recordati", Inches(1.8), Inches(6.9),
      Inches(6), Inches(0.4), size=9, color=GREY)
slide_number(s1, 1)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — Le Problème
# ═══════════════════════════════════════════════════════════════
s2 = add_slide()
txbox(s2, "Le Problème", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s2, y=Inches(1.0), w=Inches(2))
txbox(s2, "Qu'est-ce qui est cassé aujourd'hui ?",
      Inches(0.5), Inches(1.1), Inches(10), Inches(0.5),
      size=14, italic=True, color=GREY)

cx = [Inches(0.5), Inches(4.6), Inches(8.7)]
titles = ["⏳  Décisions Ralenties", "🧩  Systèmes Déconnectés", "❓  Traçabilité Absente"]
bullets = [
    ["Données éparpillées", "entre plusieurs sources", "— aucune vue centralisée"],
    ["FDA / NIH / ICH non liés", "Recherche manuelle", "chronophage et coûteuse"],
    ["Données réelles vs estimées", "ICH indiscernables", "— risque de non-conformité"]
]
for i in range(3):
    card(s2, cx[i], Inches(1.9), Inches(4.0), Inches(4.0),
         titles[i], bullets[i])

slide_number(s2, 2)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — La Solution
# ═══════════════════════════════════════════════════════════════
s3 = add_slide()
txbox(s3, "La Solution", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s3, y=Inches(1.0), w=Inches(1.8))
txbox(s3, "Découvrez PHARMA-INTEL v3.0",
      Inches(0.5), Inches(1.1), Inches(10), Inches(0.5),
      size=14, italic=True, color=GREY)

# Formula boxes
steps = ["Données\nFragmentées", "PHARMA-INTEL\nv3.0", "Décisions\nStratégiques"]
colors= [GREY, RED, WHITE]
bolds = [False, True, True]
for i, (label, col, bd) in enumerate(zip(steps, colors, bolds)):
    bx = Inches(0.5 + i*4.2)
    box = s3.shapes.add_shape(1, bx, Inches(1.9), Inches(3.8), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = DARKCARD if i != 1 else RGBColor(0x3D, 0x07, 0x09)
    box.line.color.rgb = RED
    txbox(s3, label, bx+Inches(0.1), Inches(1.95), Inches(3.6), Inches(1.3),
          size=16, bold=bd, color=col, align=PP_ALIGN.CENTER)
    if i < 2:
        txbox(s3, "►", Inches(4.2 + i*4.2), Inches(2.3), Inches(0.5), Inches(0.6),
              size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 3 keywords
kws = ["Centralisé", "Tracé", "Actionnable"]
for i, kw in enumerate(kws):
    txbox(s3, kw, Inches(1.0 + i*4.0), Inches(3.6), Inches(3.5), Inches(0.6),
          size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

txbox(s3, "→ Insérer ici une capture d'écran du dashboard PHARMA-INTEL v3.0",
      Inches(0.5), Inches(4.4), Inches(12), Inches(2.6),
      size=13, italic=True, color=GREY, align=PP_ALIGN.CENTER)
# Placeholder box for screenshot
ph = s3.shapes.add_shape(1, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.6))
ph.fill.solid(); ph.fill.fore_color.rgb = DARKCARD
ph.line.color.rgb = RED

slide_number(s3, 3)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — IA & Innovation
# ═══════════════════════════════════════════════════════════════
s4 = add_slide()
txbox(s4, "IA & Innovation", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s4, y=Inches(1.0), w=Inches(2.4))
txbox(s4, "L'innovation, c'est la valeur créée — pas la technologie utilisée.",
      Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
      size=13, italic=True, color=GREY)

# 3-column flow
cols_titles = ["Sources Brutes", "Moteur IA Asynchrone", "Output Décisionnel"]
cols_content = [
    ["FDA / NIH\n(données certifiées)", "Normes ICH\n(estimations)"],
    ["Classification auto\ndes données", "13 000 entités\ningérées en temps réel"],
    ["✅  Données Réelles", "⚡  Données Estimées ICH"]
]
cols_colors = [GREY, RED, WHITE]
for i in range(3):
    bx = Inches(0.5 + i*4.27)
    card(s4, bx, Inches(1.8), Inches(4.0), Inches(4.5),
         cols_titles[i], cols_content[i], title_color=cols_colors[i])
    if i < 2:
        txbox(s4, "►", Inches(4.4 + i*4.27), Inches(3.8), Inches(0.5), Inches(0.6),
              size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)

txbox(s4, "💡  Différenciateur clé : séparation automatique Réel (FDA/NIH) vs Estimé (ICH)",
      Inches(0.5), Inches(6.6), Inches(12), Inches(0.6),
      size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
slide_number(s4, 4)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — Impact & Durabilité
# ═══════════════════════════════════════════════════════════════
s5 = add_slide()
txbox(s5, "Impact & Durabilité", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s5, y=Inches(1.0), w=Inches(2.8))
txbox(s5, "Des résultats mesurables — le critère le plus important pour le jury.",
      Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
      size=13, italic=True, color=GREY)

impact_titles  = ["💼  Impact Business", "🏥  Impact Social", "🛡️  Impact Conformité"]
impact_bullets = [
    ["Décisions accélérées", "Gain de productivité", "Réduction des coûts ops"],
    ["Meilleure qualité ttt", "Time-to-Market réduit", "Innovation R&D soutenue"],
    ["Traçabilité totale BPF/ICH", "Audit trail automatisé", "Zéro risque non-conformité"]
]
for i in range(3):
    card(s5, Inches(0.5 + i*4.27), Inches(1.9), Inches(4.0), Inches(4.7),
         impact_titles[i], impact_bullets[i])

txbox(s5, "→ Quantifiez si possible : ex. \"40% de temps économisé sur l'analyse des rapports\"",
      Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
      size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)
slide_number(s5, 5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — Modèle Économique
# ═══════════════════════════════════════════════════════════════
s6 = add_slide()
txbox(s6, "Modèle Économique", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s6, y=Inches(1.0), w=Inches(2.8))
txbox(s6, "Comment la solution crée et maintient de la valeur dans le temps.",
      Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
      size=13, italic=True, color=GREY)

biz_titles  = ["🔻  RÉDUIRE", "🔼  ACCÉLÉRER", "🛡️  PROTÉGER"]
biz_bullets = [
    ["Coûts opérationnels", "Revue manuelle éliminée", "Recherche automatisée"],
    ["Prise de décision", "Time-to-Market", "Cycles d'innovation R&D"],
    ["Conformité BPF", "Audits ICH", "Risques réglementaires"]
]
for i in range(3):
    card(s6, Inches(0.5 + i*4.27), Inches(1.9), Inches(4.0), Inches(4.5),
         biz_titles[i], biz_bullets[i])

slide_number(s6, 6)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — Roadmap
# ═══════════════════════════════════════════════════════════════
s7 = add_slide()
txbox(s7, "Feuille de Route", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s7, y=Inches(1.0), w=Inches(2.4))
txbox(s7, "De la conception au déploiement — un plan concret.",
      Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
      size=13, italic=True, color=GREY)

# Timeline line
tl = s7.shapes.add_shape(1, Inches(0.5), Inches(4.1), Inches(12.3), Pt(3))
tl.fill.solid(); tl.fill.fore_color.rgb = RED; tl.line.fill.background()

phases = [
    ("✅", "Phase 1", "Architecture\n& Design UX/UI", "TERMINÉ", RGBColor(0x28,0xA7,0x45)),
    ("🔄", "Phase 2", "Tests de charge\n& Audit Sécurité", "EN COURS", RED),
    ("⏳", "Phase 3", "Déploiement\npilote interne", "T3 2025", GREY),
    ("⏳", "Phase 4", "Mise à l'échelle\n& Export certifié", "T4 2025", GREY),
]
for i, (icon, ph, desc, status, scol) in enumerate(phases):
    bx = Inches(0.4 + i*3.23)
    # dot on timeline
    dot = s7.shapes.add_shape(9, bx+Inches(1.35), Inches(3.88), Inches(0.44), Inches(0.44))
    dot.fill.solid(); dot.fill.fore_color.rgb = scol; dot.line.fill.background()

    txbox(s7, icon + "  " + ph, bx, Inches(2.0), Inches(3.1), Inches(0.55),
          size=13, bold=True, color=scol, align=PP_ALIGN.CENTER)
    txbox(s7, desc, bx, Inches(2.55), Inches(3.1), Inches(1.1),
          size=11, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s7, status, bx, Inches(4.4), Inches(3.1), Inches(0.5),
          size=11, bold=True, color=scol, align=PP_ALIGN.CENTER)

slide_number(s7, 7)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — L'Équipe
# ═══════════════════════════════════════════════════════════════
s8 = add_slide()
txbox(s8, "L'Équipe", Inches(0.5), Inches(0.3), Inches(10), Inches(0.7),
      size=32, bold=True, color=WHITE)
red_bar(s8, y=Inches(1.0), w=Inches(1.4))
txbox(s8, "Les architectes du projet — des expertises complémentaires au service de l'innovation.",
      Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
      size=13, italic=True, color=GREY)

roles = [
    ("👤", "Membre 1", "Architecture Technique\n& Full-Stack"),
    ("👤", "Membre 2", "Design UX/UI\n& Charte Opalia"),
    ("👤", "Membre 3", "Data & Conformité ICH"),
    ("👤", "Membre 4", "Expertise Pharmaco\n& Domaine Métier"),
]
for i, (icon, name, role) in enumerate(roles):
    bx = Inches(0.3 + i*3.25)
    box = s8.shapes.add_shape(1, bx, Inches(2.0), Inches(3.1), Inches(4.5))
    box.fill.solid(); box.fill.fore_color.rgb = DARKCARD; box.line.color.rgb = RED

    txbox(s8, icon, bx+Inches(0.05), Inches(2.15), Inches(3.0), Inches(0.8),
          size=40, align=PP_ALIGN.CENTER)
    txbox(s8, name, bx+Inches(0.05), Inches(3.15), Inches(3.0), Inches(0.5),
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s8, role, bx+Inches(0.05), Inches(3.65), Inches(3.0), Inches(0.9),
          size=11, color=RED, align=PP_ALIGN.CENTER)

txbox(s8, "→ Remplacez les icônes 👤 par vos photos et les noms par les noms réels de l'équipe.",
      Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
      size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)
slide_number(s8, 8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — Call to Action
# ═══════════════════════════════════════════════════════════════
s9 = add_slide()

# Ambient red glow (bottom-left shape)
glow = s9.shapes.add_shape(9, Inches(-1), Inches(4.5), Inches(5), Inches(5))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x30, 0x04, 0x06)
glow.line.fill.background()

txbox(s9, "Façonnons l'avenir d'Opalia Recordati",
      Inches(1.0), Inches(0.8), Inches(11), Inches(1.4),
      size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

sep9 = s9.shapes.add_shape(1, Inches(4.5), Inches(2.3), Inches(4.3), Pt(2))
sep9.fill.solid(); sep9.fill.fore_color.rgb = RED; sep9.line.fill.background()

txbox(s9,
      "En transformant la donnée pharmaceutique brute en intelligence décisionnelle,\n"
      "PHARMA-INTEL v3.0 accélère l'innovation, renforce la conformité,\n"
      "et façonne l'avenir de la santé propulsée par l'IA.\n\n"
      "La plateforme est prête. La valeur est démontrée.",
      Inches(1.0), Inches(2.5), Inches(11), Inches(2.0),
      size=14, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# Big red CTA button
btn = s9.shapes.add_shape(1, Inches(4.0), Inches(4.8), Inches(5.3), Inches(0.85))
btn.fill.solid(); btn.fill.fore_color.rgb = RED; btn.line.fill.background()
txbox(s9, "ACCÉDER À LA PLATEFORME  →",
      Inches(4.0), Inches(4.82), Inches(5.3), Inches(0.8),
      size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s9, "« L'intelligence au service de la molécule. »",
      Inches(1.0), Inches(5.9), Inches(11), Inches(0.6),
      size=16, italic=True, color=GREY, align=PP_ALIGN.CENTER)

txbox(s9, "© 2025 Opalia Recordati  —  Tous droits réservés",
      Inches(1.0), Inches(7.0), Inches(11), Inches(0.35),
      size=9, color=GREY, align=PP_ALIGN.CENTER)
slide_number(s9, 9)

# ── Save ──────────────────────────────────────────────────────
out = "PHARMA_INTEL_v3_Pitch.pptx"
prs.save(out)
print("Fichier genere avec succes : " + out)
